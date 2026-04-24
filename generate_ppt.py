from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    # Create a new presentation object
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "DataGuard AI"
    subtitle.text = "AI-powered dataset health analysis for Data Scientists\nRight inside VS Code & Antigravity"
    
    # Speaker Notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Welcome everyone. Today I'm excited to present DataGuard AI, a new extension for VS Code and Antigravity designed to bring AI-powered dataset health analysis directly into the workflow of data scientists and engineers."

    # --- Slide 2: The Problem ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "The Data Quality Challenge"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    tf.text = "Data scientists spend up to 80% of their time cleaning and preparing data."
    
    p = tf.add_paragraph()
    p.text = "Hidden issues in datasets lead to faulty models:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Missing values and nulls"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Outliers and extreme values"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Class imbalances in categorical data"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Inconsistent data types"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Traditional EDA (Exploratory Data Analysis) requires writing boilerplate code (e.g., df.info(), df.describe(), df.isnull().sum()) every time."
    p.level = 0
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "As we all know, data cleaning is the most time-consuming part of data science. Hidden issues like missing values, outliers, or class imbalances can silently ruin a model's performance. Traditionally, we have to write the same pandas boilerplate code over and over to inspect our data. This interrupts the flow and takes up valuable time."

    # --- Slide 3: Introducing DataGuard AI ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "What is DataGuard AI?"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "DataGuard AI integrates directly into your editor to provide instant, automated data profiling."
    
    p = tf.add_paragraph()
    p.text = "Zero Boilerplate: No need to write pandas profiling scripts."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Context-Aware: Appears as a CodeLens above pandas/polars read functions (e.g., pd.read_csv)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Auto-Trigger: Analyzes global CSV/Parquet/JSON files as soon as you open them."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI-Powered: Uses Google Gemini AI to provide actionable summaries and advice."
    p.level = 1
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "DataGuard AI changes this by bringing automated profiling directly into your IDE. It requires zero boilerplate. It intelligently detects when you are loading data in Python and offers a CodeLens button. Better yet, it automatically analyzes any CSV, Parquet, or JSON file you open anywhere on your system. And it leverages Google Gemini to provide human-readable, actionable advice."

    # --- Slide 4: Key Features ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Core Features & Dashboard"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Rich Interactive Webview Dashboard"
    
    p = tf.add_paragraph()
    p.text = "Health Score Gauge (0–100)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Missing Values Bar Chart"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Types Donut Chart"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Per-column statistics & Outlier badges"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Inline Editor Decorations"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Red wavy underlines directly on the data-loading code if the health score is low."
    p.level = 1
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "The extension provides a beautiful, interactive dashboard built with Tailwind CSS and Chart.js. It gives you a quick visual overview: a health score gauge, missing value charts, and data type distributions. It also highlights specific column details. Furthermore, it adds inline editor decorations—red wavy underlines—right on your Python code if the loaded dataset is unhealthy, alerting you immediately."

    # --- Slide 4.5: Analysis Capabilities ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Comprehensive Analysis Capabilities"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "DataGuard AI covers 6 key pillars of dataset health:"
    
    p = tf.add_paragraph()
    p.text = "Missing Values: Visualizing gaps in dataset entries"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Outlier Detection: Identifying anomalies in data distribution"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Class Imbalance: Highlighting unequal class distributions"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Types: Categorizing information in datasets"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Data Processing: Visualizing the flow of data analysis"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Analysis Techniques: Integrating multiple data evaluation methods"
    p.level = 1
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Our platform offers a complete suite of analysis capabilities, spanning missing values, outlier detection, class imbalance, and data types, all the way to advanced data processing tracking and integrated analysis techniques."

    # --- Slide 5: Health Score Calculation ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "How The Health Score Works"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "A simple, standardized metric for data quality (0-100)"
    
    p = tf.add_paragraph()
    p.text = "Formula: score = 100 - (avg_missing_pct × 0.5) - (outlier_column_count × 5)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Categories:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "🟢 80–100: Healthy dataset"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "🟡 50–79: Moderate issues (e.g., some missing data, few outliers)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "🔴 0–49: Critical data quality problems"
    p.level = 1
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "We wanted a quick, recognizable metric, so we built the Health Score. It starts at 100 and penalizes the score based on the average percentage of missing values and the number of columns containing significant outliers. 80 to 100 means you are good to go, 50 to 79 means you should investigate, and below 50 means critical cleaning is required."

    # --- Slide 6: Architecture under the hood ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Technical Architecture"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "A hybrid Extension & Sidecar model"
    
    p = tf.add_paragraph()
    p.text = "Extension Host (TypeScript): Handles VS Code APIs, UI, global file watching, and CodeLens."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Analysis Engine (Python): A lightweight sidecar process using Pandas to perform the actual heavy lifting."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Cross-Platform Support: Uses dynamic Python path detection to work seamlessly on Windows, macOS, and Linux."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Gemini Integration: Generates human-readable summaries locally via the sidecar."
    p.level = 1
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Technically, DataGuard AI uses a hybrid architecture. The main extension is written in TypeScript, managing the editor integration and UI. When an analysis is triggered, it spawns a lightweight Python sidecar process. This Python engine uses pandas and numpy to rapidly compute statistics and calls the Google Gemini API for the intelligent summary, returning the JSON payload back to the UI."

    # --- Slide 7: Availability & Next Steps ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Availability"
    
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Now available for download!"
    
    p = tf.add_paragraph()
    p.text = "Version: v0.1.0"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Compatibility: VS Code (1.85+) and Antigravity"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Open Source: Fully open source under the MIT License."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Next Steps:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Add automated data cleaning suggestions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Support for database connections (SQL)."
    p.level = 1
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "DataGuard AI version 0.1.0 is open source, MIT licensed, and fully compatible with both VS Code and Antigravity. In the future, we plan to add features like one-click automated cleaning code generation and direct database profiling. Thank you for your time, and I encourage you to try it out on your datasets."

    # --- Slide 8: Q&A ---
    slide = prs.slides.add_slide(title_only_layout)
    title = slide.shapes.title
    title.text = "Thank You! \nQuestions?"
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Thank you. I will now take any questions you might have."

    # Save the presentation
    prs.save('DataGuard_AI_Presentation.pptx')
    print("Presentation saved as 'DataGuard_AI_Presentation.pptx'")

if __name__ == '__main__':
    create_presentation()
